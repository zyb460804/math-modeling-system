#!/usr/bin/env python3
"""build_defense_pptx.py — 国赛中文答辩 PPT 生成器（defense-ppt-builder-zh skill）.

从论文源稿 + 结果 + 图表直接生成 .pptx（非大纲）。设计参数遵循 SKILL.md：
16:9 / 微软雅黑 / 白底 + 深蓝标题(#1F4E79) + 深灰正文(#333333) + 红色强调(#C00000) / 页码右下。

用法（与 SKILL.md Step 4 一致）：
    python .claude/skills/defense-ppt-builder-zh/scripts/build_defense_pptx.py \
        --source paper_output/final_paper_source.md \
        --figures paper_output/figures/ \
        --results paper_output/results/model_results.json \
        --qa-bank paper_output/qa/defense_qa_bank.md \
        --topic-type C \
        --output paper_output/defense.pptx

缺输入文件不阻塞：对应内容降级为 [待补] 占位，脚本退出码仍为 0（生成是产出，
完整性由 SKILL.md Step 5 自审循环把关）。
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TITLE_BLUE = RGBColor(0x1F, 0x4E, 0x79)
BODY_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_RED = RGBColor(0xC0, 0x00, 0x00)
FONT = "微软雅黑"
NUM_FONT = "Times New Roman"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

TOPIC_PAGES = {  # 题型 → 推荐“每问”页组织（见 SKILL.md Step 2）
    "A": ("机理推导 + 数值验证 + 灵敏度", "推导 40% / 结果 40% / 总结 20%"),
    "B": ("模型构建 + 求解策略 + 方案对比", "建模 30% / 求解 40% / 结果 30%"),
    "C": ("数据处理 + 评价指标 + 排名结果", "数据 30% / 建模 40% / 结果 30%"),
    "D": ("数据特征 + 预测模型 + 精度验证", "数据 25% / 建模 45% / 检验 30%"),
    "E": ("问题定义 + 方案设计 + 可行性", "定义 30% / 方案 50% / 验证 20%"),
    "F": ("问题定义 + 方案设计 + 可行性", "定义 30% / 方案 50% / 验证 20%"),
}


def read_text(p: Path | None) -> str:
    if p and p.exists():
        return p.read_text(encoding="utf-8", errors="replace")
    return ""


def load_results(p: Path | None) -> dict:
    if p and p.exists():
        try:
            return json.loads(p.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            return {}
    return {}


def parse_sections(md_text: str) -> dict[str, str]:
    """按二级/三级标题切块，取每块前几行有效文字作要点。"""
    sections: dict[str, str] = {}
    cur, buf = "", []
    for line in md_text.splitlines():
        m = re.match(r"^#{1,3}\s+(.*)", line)
        if m:
            if cur:
                sections[cur] = "\n".join(buf).strip()
            cur, buf = m.group(1).strip(), []
        else:
            buf.append(line)
    if cur:
        sections[cur] = "\n".join(buf).strip()
    return sections


def pick(sections: dict[str, str], *keywords: str, limit: int = 4) -> list[str]:
    """取首个命中关键词的标题块，提炼为 ≤limit 条短要点（每条 ≤40 字）。"""
    for kw_list in keywords:
        for title, body in sections.items():
            if all(k in title for k in kw_list):
                pts = [re.sub(r"[*#`>|]", "", l).strip() for l in body.splitlines()]
                pts = [p for p in pts if len(p) > 4][:limit]
                if pts:
                    return [p[:40] + ("…" if len(p) > 40 else "") for p in pts]
    return []


def list_figures(fig_dir: Path | None) -> list[Path]:
    if not fig_dir or not fig_dir.exists():
        return []
    return sorted([f for f in fig_dir.iterdir()
                   if f.suffix.lower() in {".png", ".jpg", ".jpeg"}])


def add_slide(prs: Presentation, total: list[int]) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    total[0] += 1
    box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4))
    tf = box.text_frame
    tf.text = f"{total[0]}"
    tf.paragraphs[0].runs[0].font.size = Pt(12)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide


def put_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(32), True, TITLE_BLUE, FONT


def put_bullets(slide, points: list[str], top=Inches(1.5), height=Inches(5.0), left=None):
    if not points:
        points = ["[待补]"]
    if left is None:
        box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12.0), height)
    else:
        box = slide.shapes.add_textbox(left, top, Inches(5.6), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, p in enumerate(points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = f"• {p}"
        run.font.size, run.font.color.rgb, run.font.name = Pt(20), BODY_GRAY, FONT


def put_figure(slide, fig: Path | None, left=Inches(0.6), top=Inches(1.5),
               width=Inches(6.5), height=Inches(5.2)):
    if fig and fig.exists():
        slide.shapes.add_picture(str(fig), left, top, width=width, height=height)
    else:
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.text = "[待补图表]"
        tf.paragraphs[0].runs[0].font.size = Pt(24)
        tf.paragraphs[0].runs[0].font.color.rgb = ACCENT_RED


def key_numbers(results: dict, limit: int = 4) -> list[str]:
    """从 model_results.json 顶层标量提炼关键数字证据。"""
    out = []
    for k, v in results.items():
        if isinstance(v, (int, float)):
            out.append(f"{k} = {v:g}")
        elif isinstance(v, dict):
            for k2, v2 in list(v.items())[:2]:
                if isinstance(v2, (int, float)):
                    out.append(f"{k}.{k2} = {v2:g}")
        if len(out) >= limit:
            break
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="国赛中文答辩 PPT 生成器")
    ap.add_argument("--source", type=Path, default=Path("paper_output/final_paper_source.md"))
    ap.add_argument("--figures", type=Path, default=Path("paper_output/figures/"))
    ap.add_argument("--results", type=Path, default=Path("paper_output/results/model_results.json"))
    ap.add_argument("--qa-bank", type=Path, default=Path("paper_output/qa/defense_qa_bank.md"))
    ap.add_argument("--topic-type", default="C", choices=list("ABCDEF"))
    ap.add_argument("--output", type=Path, default=Path("paper_output/defense.pptx"))
    args = ap.parse_args()

    md = read_text(args.source)
    if not md:
        print(f"[warn] 论文源稿缺失或为空: {args.source} —— 页面将出现 [待补]")
    sections = parse_sections(md)
    results = load_results(args.results)
    figs = list_figures(args.figures)
    fig_i = 0
    focus, timing = TOPIC_PAGES[args.topic_type.upper()]

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    total = [0]

    # P1 封面
    s = add_slide(prs, total)
    title = (re.search(r"^#\s+(.*)", md).group(1).strip()
             if re.search(r"^#\s+(.*)", md) else f"{args.topic_type} 题 [待补题目]")
    box = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.4))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].runs[0]
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(40), True, TITLE_BLUE, FONT
    box2 = s.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.3), Inches(0.8))
    tf2 = box2.text_frame
    tf2.text = "队号：[待补]    成员：[待补]"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.paragraphs[0].runs[0].font.size = Pt(20)
    tf2.paragraphs[0].runs[0].font.color.rgb = BODY_GRAY

    # P2 问题概述 + P3 总体思路（模型框架图 = 全场最重要的一页）
    s = add_slide(prs, total)
    put_title(s, "问题概述")
    put_bullets(s, pick(sections, ("问题", "重述"), ("摘要",)) or ["[待补] 一句话讲清题目本质"])
    s = add_slide(prs, total)
    put_title(s, "总体思路")
    put_bullets(s, pick(sections, ("思路",), ("模型", "建立"), ("方法",)) or ["[待补] 技术路线"],
                top=Inches(1.5), height=Inches(2.2))
    put_figure(s, figs[fig_i] if fig_i < len(figs) else None,
               left=Inches(0.6), top=Inches(3.8), width=Inches(12.1), height=Inches(3.2))
    fig_i += 1

    # 每问 3 页：建模 / 结果 / 检验
    q_titles = [t for t in sections if re.match(r"^问题[一二三四五六1-6]", t)]
    if not q_titles:
        q_titles = ["问题一", "问题二", "问题三"]
    for qt in q_titles[:4]:
        core = re.sub(r"^问题[一二三四五六1-6][：:\s]*", "", qt) or "[待补核心模型]"
        for kind in ("模型建立", "求解与结果", "检验"):
            s = add_slide(prs, total)
            put_title(s, f"{qt.split('：')[0]}：{kind}")
            pts = pick(sections, (core[:4],)) if kind == "模型建立" else []
            if kind == "求解与结果":
                pts = key_numbers(results) or ["[待补关键数字]"]
            if kind == "检验":
                pts = pick(sections, ("灵敏",), ("检验",), ("误差",)) or ["[待补检验证据]"]
            put_bullets(s, pts, top=Inches(1.5), height=Inches(2.0))
            put_figure(s, figs[fig_i] if fig_i < len(figs) else None,
                       left=Inches(0.6), top=Inches(3.6), width=Inches(12.1), height=Inches(3.4))
            fig_i += 1

    # 创新点（数字说话）/ 模型评价 / 问答预备页
    s = add_slide(prs, total)
    put_title(s, "创新点总结")
    put_bullets(s, [f"创新点{i+1}：{n}" for i, n in enumerate(key_numbers(results, 3))]
                or ["[待补] 3 条创新点，每条 1 句话 + 1 个数字证据"])
    s = add_slide(prs, total)
    put_title(s, "模型评价")
    put_bullets(s, ["优点：[待补]", "局限：[待补]", "改进方向：[待补]"])
    s = add_slide(prs, total)
    put_title(s, "问答预备")
    put_figure(s, figs[-1] if figs else None,
               left=Inches(0.6), top=Inches(1.5), width=Inches(12.1), height=Inches(5.4))

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"[ok] 已生成 {args.output}（{total[0]} 页，题型 {args.topic_type.upper()}：{focus}｜{timing}）")
    print("[next] 按 SKILL.md Step 5 自审循环检查：文字溢出 / 图表 ≥150DPI / 图表占比 ≥50%，最多 3 轮修正")
    return 0


if __name__ == "__main__":
    sys.exit(main())
